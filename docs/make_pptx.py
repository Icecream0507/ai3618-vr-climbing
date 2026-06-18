#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Generate Summit VR coursework presentation as a .pptx (SJTU red/white theme).
Run: python docs/make_pptx.py  ->  docs/SummitVR_Pre.pptx
Academic project-report register (AI3618 课程作业汇报).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
DEMO = os.path.join(HERE, "..", "Demo")
TH = os.path.join(HERE, "_thumbs")

# SJTU palette
RED   = RGBColor(0xC8, 0x10, 0x2E)
REDD  = RGBColor(0x9E, 0x0B, 0x22)
REDL  = RGBColor(0xE8, 0x34, 0x4C)
INK   = RGBColor(0x1C, 0x1F, 0x26)
MUT   = RGBColor(0x6B, 0x72, 0x80)
PAPER2= RGBColor(0xFB, 0xF7, 0xF8)
LINE  = RGBColor(0xE7, 0xD6, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
REDWASH = RGBColor(0xFB, 0xE9, 0xEC)
GREEN = RGBColor(0x1F, 0xAA, 0x59)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def fill(sp, color):
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()

def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h); fill(sp, color); return sp

def grad_band(s, x, y, w, h):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.line.fill.background()
    f = sp.fill; f.gradient()
    f.gradient_stops[0].color.rgb = RED
    f.gradient_stops[1].color.rgb = REDL
    try: f.gradient_angle = 0.0
    except Exception: pass
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line=1.12):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs[0], tuple): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line
        for (t, sz, col, b) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = b
            r.font.name = "Microsoft YaHei"
    return tb

def kicker(s, t, y=Inches(0.42)):
    rect(s, Inches(0.55), y+Inches(0.02), Inches(0.06), Inches(0.30), RED)
    txt(s, Inches(0.72), y, Inches(10), Inches(0.4), [[(t, 14, RED, True)]])

def header_band(s):
    grad_band(s, 0, 0, SW, Inches(0.16))

def footer(s, left, right="上海交通大学 · AI3618"):
    rect(s, 0, SH-Inches(0.42), SW, Inches(0.42), RED)
    txt(s, Inches(0.4), SH-Inches(0.40), Inches(8), Inches(0.36),
        [[(left, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, SW-Inches(5.4), SH-Inches(0.40), Inches(5), Inches(0.36),
        [[(right, 11, WHITE, False)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def card(s, x, y, w, h, red_bg=False):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if red_bg:
        f = sp.fill; f.gradient()
        f.gradient_stops[0].color.rgb = RED
        f.gradient_stops[1].color.rgb = REDD
        sp.line.fill.background()
    else:
        fill(sp, PAPER2); sp.line.color.rgb = LINE; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def set_run_font(r, name="Segoe UI Emoji"):
    r.font.name = name

def bullets(s, x, y, w, h, items, size=15, gap=7):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.25
        for (t, b, c) in it:
            r = p.add_run(); r.text = t; r.font.size = Pt(size)
            r.font.bold = b; r.font.color.rgb = c; r.font.name = "Microsoft YaHei"
    return tb

def add_video(s, name, x, y, w, h):
    mp4 = os.path.join(DEMO, "SummitVR_%s.mp4" % name)
    thumb = os.path.join(TH, "%s.png" % name)
    poster = thumb if os.path.exists(thumb) else None
    try:
        return s.shapes.add_movie(mp4, x, y, w, h, poster_frame_image=poster,
                                  mime_type="video/mp4")
    except Exception as e:
        print("  movie embed failed (%s): %s -> image" % (name, e))
        return s.shapes.add_picture(poster, x, y, w, h) if poster else None

def codebox(s, x, y, w, h, lines, size=11):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill(sp, PAPER2); sp.line.color.rgb = LINE; sp.line.width = Pt(1); sp.shadow.inherit = False
    rect(s, x, y, Inches(0.08), h, RED)
    tb = s.shapes.add_textbox(x+Inches(0.22), y+Inches(0.12), w-Inches(0.4), h-Inches(0.24))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (t, hot) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.18
        r = p.add_run(); r.text = t; r.font.size = Pt(size)
        r.font.name = "Consolas"; r.font.color.rgb = RED if hot else RGBColor(0x33,0x37,0x3F)
        r.font.bold = hot

def legend(s, x, y):
    items = [("手","F5B301"),("脚","FF7A1A"),("任意","9B51E0"),("登顶","1FAA59"),("易碎","E23B3B"),("休息","2D8CFF")]
    cx = x
    for label, hexc in items:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, y, Inches(0.16), Inches(0.16))
        fill(d, RGBColor.from_string(hexc)); d.shadow.inherit=False
        txt(s, cx+Inches(0.2), y-Inches(0.05), Inches(0.7), Inches(0.3), [[(label, 11, MUT, False)]])
        cx = cx + Inches(0.2) + Inches(0.14*len(label)+0.45)

def stick_climber(s, ox, oy, scale, both_right=False, footed=False):
    U = scale; col = INK
    def line(x1, y1, x2, y2, c=col, wpt=2.2):
        ln = s.shapes.add_connector(2, ox+x1, oy+y1, ox+x2, oy+y2)
        ln.line.color.rgb = c; ln.line.width = Pt(wpt); return ln
    def ball(cx, cy, r, c):
        b = s.shapes.add_shape(MSO_SHAPE.OVAL, ox+cx-r, oy+cy-r, 2*r, 2*r)
        fill(b, c); b.shadow.inherit = False; return b
    line(U*1.0, U*0.9, U*1.0, U*2.4)
    ball(U*1.0, U*0.6, U*0.32, col)
    if both_right:
        line(U*1.0, U*1.2, U*1.9, U*0.7); line(U*1.0, U*1.2, U*1.7, U*1.1)
    else:
        line(U*1.0, U*1.2, U*1.8, U*0.8); line(U*1.0, U*1.2, U*0.2, U*0.8)
    line(U*1.0, U*2.4, U*1.5, U*3.2)
    if footed:
        line(U*1.0, U*2.4, U*0.2, U*3.0); ball(U*0.2, U*3.0, U*0.18, RGBColor(0xFF,0x7A,0x1A))
    else:
        line(U*1.0, U*2.4, U*0.7, U*3.2)

print("helpers ready")
# <<<SLIDES>>>

# ============== 0 — COVER ==============
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, Inches(8.4), 0, SW-Inches(8.4), SH, RED)
grad_band(s, 0, 0, Inches(8.4), Inches(0.18))
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.4), Inches(5.4), Inches(2.2), Inches(2.1))
fill(tri, REDD); tri.shadow.inherit = False; tri.rotation = 180
txt(s, Inches(0.8), Inches(1.0), Inches(7.4), Inches(0.5),
    [[("上海交通大学 · SHANGHAI JIAO TONG UNIVERSITY", 14, MUT, True)]])
txt(s, Inches(0.75), Inches(1.9), Inches(7.6), Inches(1.4),
    [[("Summit ", 70, INK, True), ("VR", 70, RED, True)]])
rect(s, Inches(0.85), Inches(3.35), Inches(2.2), Pt(4), RED)
txt(s, Inches(0.8), Inches(3.6), Inches(7.4), Inches(1.2),
    [[("面向纯手柄 VR 攀岩的平衡与脚法机制", 23, INK, True)],
     [("课程项目 · 设计、实现与评测方案", 14, MUT, False)]],
    space_after=8, line=1.3)
txt(s, Inches(0.8), Inches(5.25), Inches(7.4), Inches(0.5),
    [[("AI3618《虚拟现实技术》· 课程作业汇报", 15, REDD, True)]])
txt(s, Inches(0.8), Inches(6.4), Inches(7.4), Inches(0.5),
    [[("第 N 组　薛俊智 · 吴一轩 · 陶锐 · 邹沛霖 · 孙艺豪", 13, MUT, False)]])
txt(s, Inches(8.8), Inches(2.6), Inches(4.2), Inches(3.0),
    [[("研究问题", 16, RGBColor(0xFF,0xD9,0xDF), True)],
     [("", 6, WHITE, False)],
     [("如何在仅有头显与两个", 17, WHITE, False)],
     [("手柄的消费级 VR 上，", 17, WHITE, False)],
     [("重建真实攀岩中的", 17, WHITE, False)],
     [("身体平衡与脚法？", 19, WHITE, True)]],
    space_after=3, line=1.3)
print("cover done")

# ============== 1 — TEAM / 分工 ==============
s = slide(); header_band(s)
kicker(s, "团队与分工")
txt(s, Inches(0.55), Inches(0.85), Inches(12), Inches(0.7),
    [[("第 N 组 · 五人分工", 30, INK, True)]])
txt(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
    [[("课程要求每人主负责一个模块；平衡与脚法机制为全组共同设计。", 13, MUT, False)]])
team = [("P1","薛俊智","攀爬系统","反向位移移动 · 抓握 · 平衡与脚法核心算法"),
        ("P2","吴一轩","玩法规则","状态机 · 计时与跌落 · 体力 · 线路设计"),
        ("P3","陶锐","场景与美术","墙体 · 握点造型 · 灯光 · 角色化身"),
        ("P4","邹沛霖","XR 集成与验证","设备接入 · 输入 · 仿真自检 · 构建"),
        ("P5","孙艺豪","UX·音频·报告","HUD · 音效 · 报告统稿 · 视频与汇报")]
ty = Inches(2.35)
for code, nm, role, tk in team:
    card(s, Inches(0.6), ty, Inches(12.1), Inches(0.82))
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), ty+Inches(0.16), Inches(0.5), Inches(0.5))
    fill(badge, RED); badge.shadow.inherit=False
    bp = badge.text_frame.paragraphs[0]; bp.alignment=PP_ALIGN.CENTER
    br = bp.add_run(); br.text=code; br.font.size=Pt(15); br.font.bold=True
    br.font.color.rgb=WHITE; br.font.name="Microsoft YaHei"
    txt(s, Inches(1.6), ty+Inches(0.18), Inches(1.6), Inches(0.5), [[(nm, 18, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.2), ty+Inches(0.18), Inches(2.4), Inches(0.5), [[(role, 14, RED, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.7), ty+Inches(0.18), Inches(6.8), Inches(0.5), [[(tk, 12.5, MUT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    ty = ty + Inches(0.92)
footer(s, "团队与分工")
print("team done")

# ============== 2 — AGENDA ==============
s = slide(); header_band(s)
kicker(s, "汇报提纲")
txt(s, Inches(0.55), Inches(0.85), Inches(12), Inches(0.7),
    [[("汇报提纲", 30, INK, True)]])
agenda = [
  ("01","研究背景与问题","VR 攀岩的现状与所缺失的能力"),
  ("02","相关工作","脚法、重心与平衡模型的已有研究"),
  ("03","方法与系统设计","三点追踪下的平衡与脚法建模"),
  ("04","核心算法","支撑区间平衡判定 · 脚法抽象"),
  ("05","实现与演示","系统实现 · 线路设计 · 运行演示"),
  ("06","验证、评测与总结","开发验证方法 · 评测方案 · 结论"),
]
colx = [Inches(0.6), Inches(6.7)]
for i, (n, h, body) in enumerate(agenda):
    cx = colx[i // 3]; cy = Inches(2.0) + Inches(1.55) * (i % 3)
    card(s, cx, cy, Inches(5.9), Inches(1.35))
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, cx+Inches(0.25), cy+Inches(0.32), Inches(0.7), Inches(0.7))
    fill(badge, RED); badge.shadow.inherit=False
    bp = badge.text_frame.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = n; br.font.size = Pt(20); br.font.bold = True
    br.font.color.rgb = WHITE; br.font.name = "Microsoft YaHei"
    txt(s, cx+Inches(1.2), cy+Inches(0.22), Inches(4.5), Inches(0.45), [[(h, 17, INK, True)]])
    txt(s, cx+Inches(1.2), cy+Inches(0.72), Inches(4.6), Inches(0.5), [[(body, 12, MUT, False)]], line=1.2)
footer(s, "汇报提纲")
print("agenda done")
# <<<NEXT>>>

# ============== 3 — 研究背景 ==============
s = slide(); header_band(s)
kicker(s, "一 · 研究背景")
txt(s, Inches(0.55), Inches(0.85), Inches(12), Inches(0.7),
    [[("VR 攀岩的交互，长期停留在「纯手」", 28, INK, True)]])
bullets(s, Inches(0.6), Inches(1.85), Inches(6.3), Inches(3.2), [
    [("现有成熟作品（如 The Climb、Gorilla Tag）以「抓握—反向位移」为核心：", False, INK)],
    [("• 玩家用双手抓点并将身体拉起，体验舒适、晕动较低；", False, INK)],
    [("• 失败仅来自两种情形——", False, INK), ("松手", True, RED), (" 或 ", False, INK), ("臂展不足够不到", True, RED), ("。", False, INK)],
    [("而真实抱石中，决定成败的往往是", False, INK), ("重心与脚的配合", True, RED),
     ("：身体若偏离支撑范围，会绕支点旋出（俗称「开门 / barn-door」）。", False, INK)],
    [("近年平面作品（New Heights, 2026）将平衡与脚法作为核心卖点，反映出此类需求。", False, MUT)],
], size=14, gap=10)
# right illustration: barn-door
c = card(s, Inches(7.1), Inches(1.85), Inches(5.6), Inches(3.6))
txt(s, Inches(7.1), Inches(2.0), Inches(5.6), Inches(0.35),
    [[("「开门」效应示意", 13, RED, True)]], align=PP_ALIGN.CENTER)
rect(s, Inches(7.5), Inches(2.5), Inches(4.8), Inches(2.6), RGBColor(0xF3,0xE7,0xE9))
stick_climber(s, Inches(8.1), Inches(2.7), Emu(int(Inches(0.62))), both_right=True, footed=False)
ar = s.shapes.add_shape(MSO_SHAPE.CURVED_RIGHT_ARROW, Inches(10.4), Inches(3.0), Inches(0.9), Inches(1.4))
fill(ar, RED); ar.shadow.inherit=False
txt(s, Inches(7.5), Inches(4.78), Inches(4.8), Inches(0.3),
    [[("双手同侧 → 重心偏出支撑 → 身体旋出脱落", 11, MUT, False)]], align=PP_ALIGN.CENTER)
footer(s, "一 · 研究背景")
print("bg done")

# ============== 4 — 问题分析 ==============
s = slide(); header_band(s)
kicker(s, "一 · 问题分析")
txt(s, Inches(0.55), Inches(0.85), Inches(12), Inches(0.7),
    [[("难点：消费级 VR ", 28, INK, True), ("看不到脚", 28, RED, True)]])
left = card(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(3.6))
txt(s, Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.4), [[("硬件约束", 15, RED, True)]])
bullets(s, Inches(0.85), Inches(2.5), Inches(5.5), Inches(2.8), [
    [("消费级头显仅提供 3 个追踪点：", False, INK), ("头显 + 两手柄", True, REDD), ("；", False, INK)],
    [("无脚部 / 躯干追踪，下半身姿态不可观测；", False, INK)],
    [("用 IK 或生成式方法反解腿部，在攀岩这类极端姿势下易失真、穿模。", False, INK)],
], size=13.5, gap=9)
right = card(s, Inches(6.8), Inches(1.85), Inches(5.9), Inches(3.6))
txt(s, Inches(7.05), Inches(2.0), Inches(5.4), Inches(0.4), [[("我们要回答的问题", 15, RED, True)]])
bullets(s, Inches(7.05), Inches(2.5), Inches(5.4), Inches(2.8), [
    [("能否在不增加任何硬件的前提下，", False, INK)],
    [("用现有的 3 个追踪点，", False, INK), ("近似", True, REDD), ("出「重心是否落在支撑范围内」？", False, INK)],
    [("并把脚法表达为「是否提供了有效支撑」，", False, INK)],
    [("从而让平衡与脚法成为", False, INK), ("真正影响成败的玩法", True, RED), ("。", False, INK)],
], size=13.5, gap=9)
footer(s, "一 · 问题分析")
print("problem done")

# ============== 5 — 相关工作 ==============
s = slide(); header_band(s)
kicker(s, "二 · 相关工作")
txt(s, Inches(0.55), Inches(0.85), Inches(12), Inches(0.7),
    [[("已有研究为方法提供了依据", 28, INK, True)]])
works = [
  ("脚的重要性","Kosmalla et al.\nCHI 2020","真实攀岩墙上的实验表明：呈现「脚」比呈现「手」更能提升动作准确感与体验。"),
  ("头作为重心","Mitsuda & Kimura\nFrontiers in VR 2026","以头显近似重心，当重心越过无支撑的脚即判定跌落；n=24 实验验证可行。"),
  ("平衡的形式化","机器人学\n(ZMP / 支撑多边形)","支撑多边形、零力矩点、capture point 给出了「重心是否在支撑区内」的理论语言。"),
]
cx = Inches(0.6)
for tag, src, body in works:
    card(s, cx, Inches(1.9), Inches(3.95), Inches(3.5))
    rect(s, cx, Inches(1.9), Inches(3.95), Inches(0.10), RED)
    txt(s, cx+Inches(0.28), Inches(2.15), Inches(3.4), Inches(0.4), [[(tag, 16, RED, True)]])
    txt(s, cx+Inches(0.28), Inches(2.7), Inches(3.4), Inches(0.7), [[(src, 12, REDD, True)]], line=1.15)
    txt(s, cx+Inches(0.28), Inches(3.55), Inches(3.45), Inches(1.7), [[(body, 12.5, INK, False)]], line=1.3)
    cx = cx + Inches(4.05)
txt(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.4),
    [[("本工作定位：", 13, REDD, True),
      ("把上述「为什么做脚」「如何用头判平衡」的结论，落实到无脚追踪的消费级硬件与可玩机制中。", 13, INK, False)]])
footer(s, "二 · 相关工作")
print("related done")
# <<<NEXT>>>

# ============== 6 — 方法概述 ==============
s = slide(); header_band(s)
kicker(s, "三 · 方法概述")
txt(s, Inches(0.55), Inches(0.85), Inches(12.4), Inches(0.7),
    [[("设计原则：", 28, INK, True), ("抽象不可观测量，只建模影响玩法的部分", 28, RED, True)]])
cards3 = [("01","🧠","以头显近似重心","头显是 3 点中唯一稳定可得的高位信号；其横向位置可近似反映躯干是否偏离支撑。"),
          ("02","🦶","以状态表达脚法","不渲染腿部几何，而将脚抽象为「踩在某握点上、提供一个支撑点」的状态。"),
          ("03","📊","以分级量表达失衡","失衡不是瞬时判定，而是带缓冲与回升的平衡量，使其可预警、可恢复。")]
cw = Inches(4.0); cx = Inches(0.6); cy = Inches(1.9)
for n, icon, h, body in cards3:
    card(s, cx, cy, cw, Inches(2.55))
    rect(s, cx, cy, cw, Inches(0.10), RED)
    ic = s.shapes.add_shape(MSO_SHAPE.OVAL, cx+Inches(0.25), cy+Inches(0.3), Inches(0.7), Inches(0.7))
    fill(ic, REDWASH); ic.line.color.rgb = RED; ic.line.width=Pt(1.2); ic.shadow.inherit=False
    ip = ic.text_frame.paragraphs[0]; ip.alignment = PP_ALIGN.CENTER
    ir = ip.add_run(); ir.text = icon; ir.font.size = Pt(22); set_run_font(ir)
    txt(s, cx+Inches(1.1), cy+Inches(0.3), cw-Inches(1.3), Inches(0.4), [[(n, 26, REDL, True)]])
    txt(s, cx+Inches(0.25), cy+Inches(1.15), cw-Inches(0.5), Inches(0.4), [[(h, 16, RED, True)]])
    txt(s, cx+Inches(0.25), cy+Inches(1.6), cw-Inches(0.5), Inches(0.9), [[(body, 12.5, INK, False)]], line=1.3)
    cx = cx + cw + Inches(0.33)
mono = card(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(0.6))
txt(s, Inches(0.85), Inches(4.8), Inches(11.6), Inches(0.4),
    [[("加法式设计：", 14, REDD, True),
      ("平衡与脚法是叠加在成熟「抓握—反向位移」基础上的可选层；若移除，系统退化为常规纯手攀岩，基础体验不受影响。", 14, INK, False)]])
footer(s, "三 · 方法概述")
print("method done")

# ============== 7 — 系统架构 ==============
s = slide(); header_band(s)
kicker(s, "三 · 系统架构")
txt(s, Inches(0.55), Inches(0.85), Inches(12.4), Inches(0.7),
    [[("模块组成与数据流", 28, INK, True)]])
arch = [
 ("输入(grip 握力)                     物理重叠检测 (Hold 层)", False),
 ("    |                                    |", False),
 ("    v                                    v", False),
 ("ClimbingHand ×2 --抓/放--> ClimbHold {角色, 类型}   (臂展约束 ~0.88m)", False),
 ("    |                                    ^", False),
 ("    |                  FootPlacementSystem --虚拟脚自动吸附到脚点", False),
 ("    v                                    |", False),
 ("ClimbController --反向位移 + 重力 + 坠落/重生--> CharacterController", False),
 ("    ^   | 接触点(手+脚)决定是否留在墙上          |", False),
 ("    |   +----------------> BalanceSystem (头=重心; 横向支撑判定)", True),
 ("    |  失衡脱落(平衡=0) <---------+", False),
 ("GameManager(状态/计时/跌落) --事件--> GameHUD(计时 / 体力条 / 平衡条)", False),
]
codebox(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(3.05), arch, size=12)
txt(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.6),
    [[("实现规模：", 13, REDD, True),
      ("14 个 C# 脚本，分为 攀爬与平衡核心 / 玩法与线路 / UI 与工具 三部分；对平衡、脚法模块的引用均做空值保护，确保可拆卸。", 13, INK, False)]], line=1.25)
footer(s, "三 · 系统架构")
print("arch done")

# ============== 8 — 平衡算法 ==============
s = slide(); header_band(s)
kicker(s, "四 · 核心算法 · 平衡判定")
txt(s, Inches(0.55), Inches(0.85), Inches(12), Inches(0.7),
    [[("将支撑面简化为", 28, INK, True), ("一维横向区间", 28, RED, True)]])
dx, dy, dw, dh = Inches(0.6), Inches(1.9), Inches(6.0), Inches(3.0)
card(s, dx, dy, dw, dh)
barY = dy + Inches(2.05)
rect(s, dx+Inches(1.2), barY, Inches(3.6), Inches(0.04), LINE)
sup = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dx+Inches(1.9), barY-Inches(0.16), Inches(2.0), Inches(0.34))
sf = sup.fill; sf.solid(); sf.fore_color.rgb = REDWASH; sup.line.color.rgb = RED; sup.line.width=Pt(1.2); sup.shadow.inherit=False
txt(s, dx+Inches(1.5), barY+Inches(0.30), Inches(3.0), Inches(0.3),
    [[("[low, high] 支撑区间", 12, RED, True)]], align=PP_ALIGN.CENTER)
def dot(cx_in, color, label, lcol):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, dx+cx_in-Inches(0.10), barY-Inches(0.10), Inches(0.20), Inches(0.20))
    fill(d, color); d.shadow.inherit=False
    txt(s, dx+cx_in-Inches(0.4), barY-Inches(0.50), Inches(0.8), Inches(0.3), [[(label, 11, lcol, True)]], align=PP_ALIGN.CENTER)
dot(Inches(2.1), RGBColor(0xF5,0xB3,0x01), "手", RGBColor(0xB5,0x86,0x0A))
dot(Inches(3.7), RGBColor(0xF5,0xB3,0x01), "手", RGBColor(0xB5,0x86,0x0A))
dot(Inches(1.9), RGBColor(0xFF,0x7A,0x1A), "脚", RGBColor(0xD8,0x59,0x0A))
head = s.shapes.add_shape(MSO_SHAPE.OVAL, dx+Inches(2.9)-Inches(0.18), dy+Inches(0.35), Inches(0.36), Inches(0.36))
fill(head, INK); head.shadow.inherit=False
rect(s, dx+Inches(2.9)-Inches(0.01), dy+Inches(0.71), Inches(0.03), Inches(1.18), INK)
txt(s, dx+Inches(2.4), dy+Inches(0.05), Inches(1.0), Inches(0.3), [[("头 = 重心", 11, INK, True)]], align=PP_ALIGN.CENTER)
txt(s, dx+Inches(0.5), dy+Inches(2.62), dw-Inches(1.0), Inches(0.3),
    [[("重心横向坐标落入区间 → 稳定", 11, MUT, False)]], align=PP_ALIGN.CENTER)
rxt = Inches(7.0)
bullets(s, rxt, Inches(1.9), Inches(5.7), Inches(2.5), [
  [("①", True, RED), (" 取所有接触点（抓握的手、踩住的脚）在墙面横轴上的投影，构成区间 [low, high]。", False, INK)],
  [("②", True, RED), (" 计算重心横向坐标到区间的有符号余量 s ∈ [−1, +1]：在内为正、在外为负。", False, INK)],
  [("③", True, RED), (" 失衡持续超过缓冲时间后，平衡量按 s 衰减；回到区间内则回升。", False, INK)],
  [("④", True, RED), (" 平衡量降至 0 触发脱落。踩脚或抓对侧点会扩大区间，使重心重新落入。", False, INK)],
], size=13, gap=8)
mty = Inches(5.5)
txt(s, Inches(0.6), mty-Inches(0.02), Inches(1.5), Inches(0.3), [[("平衡量：", 12, INK, True)]])
seg_w = Inches(1.7)
colors = [GREEN, RGBColor(0x8B,0xC3,0x4A), RGBColor(0xF5,0xB3,0x01), RGBColor(0xFF,0x7A,0x1A), RGBColor(0xE2,0x3B,0x3B)]
mx = Inches(1.9)
for i, c in enumerate(colors):
    rect(s, mx+seg_w*i, mty, seg_w, Inches(0.3), c)
txt(s, mx, mty+Inches(0.32), Inches(2.2), Inches(0.25), [[("满 = 稳定", 10, GREEN, True)]])
txt(s, mx+seg_w*4-Inches(0.6), mty+Inches(0.32), seg_w+Inches(0.7), Inches(0.25),
    [[("空 = 脱落坠落", 10, RED, True)]], align=PP_ALIGN.RIGHT)
footer(s, "四 · 平衡判定算法")
print("balance done")
# <<<NEXT>>>

# ============== 9 — 脚法建模 ==============
s = slide(); header_band(s)
kicker(s, "四 · 核心算法 · 脚法建模")
txt(s, Inches(0.55), Inches(0.85), Inches(12.4), Inches(0.7),
    [[("脚 = ", 28, INK, True), ("由握点派生的支撑状态", 28, RED, True), ("，而非腿部几何", 28, INK, True)]])
bullets(s, Inches(0.6), Inches(1.9), Inches(6.4), Inches(3.2), [
  [("吸附：", True, REDD), ("在身体下方估计落脚区，左右各选最近的可踩握点；爬过后脱离并吸附下一个（", False, INK), ("foot-gluing", True, RED), ("）。", False, INK)],
  [("作用：", True, REDD), ("每个踩住的脚向平衡判定贡献一个支撑点，从而扩大横向支撑区间。", False, INK)],
  [("为何不反解腿部？", True, REDD), ("攀岩姿势下 IK / 生成式方法易失真穿模，且与机制无关——机制只需要「脚是否提供支撑」这一信息。", False, INK)],
  [("化身仅作展示：", True, REDD), ("演示中的人形用阻尼摆 + 两段 IK 表现重量感，不参与玩法计算。", False, MUT)],
], size=14, gap=11)
# right: before/after stick comparison reused
c = card(s, Inches(7.1), Inches(1.9), Inches(5.6), Inches(3.4))
txt(s, Inches(7.1), Inches(2.05), Inches(5.6), Inches(0.35), [[("脚扩大支撑区间 → 重心回正", 13, RED, True)]], align=PP_ALIGN.CENTER)
rect(s, Inches(7.4), Inches(2.55), Inches(2.35), Inches(2.4), RGBColor(0xF3,0xE7,0xE9))
stick_climber(s, Inches(7.7), Inches(2.7), Emu(int(Inches(0.55))), both_right=True, footed=False)
txt(s, Inches(7.4), Inches(4.62), Inches(2.35), Inches(0.3), [[("无脚：偏出", 11, MUT, False)]], align=PP_ALIGN.CENTER)
rect(s, Inches(10.05), Inches(2.55), Inches(2.35), Inches(2.4), RGBColor(0xE7,0xF3,0xEB))
stick_climber(s, Inches(10.35), Inches(2.7), Emu(int(Inches(0.55))), both_right=True, footed=True)
txt(s, Inches(10.05), Inches(4.62), Inches(2.35), Inches(0.3), [[("踩脚：回正", 11, GREEN, True)]], align=PP_ALIGN.CENTER)
footer(s, "四 · 脚法建模")
print("foot done")

# ============== 10 — 关键参数 ==============
s = slide(); header_band(s)
kicker(s, "四 · 关键参数")
txt(s, Inches(0.55), Inches(0.85), Inches(12.4), Inches(0.7),
    [[("可调参数：经验取值，编辑器内可调", 28, INK, True)]])
rows = [("模块","参数","取值","含义"),
        ("BalanceSystem","supportMargin","0.06 m","支撑区间两侧的容差"),
        ("BalanceSystem","maxOvershoot","0.30 m","映射为完全失衡的偏出量"),
        ("BalanceSystem","drain / regen","0.6 / 0.7 /s","平衡量衰减 / 回升速率"),
        ("BalanceSystem","graceTime","0.35 s","开始衰减前的缓冲（抗抖动）"),
        ("FootPlacement","footReach","0.45 m","脚的吸附 / 保持半径"),
        ("ClimbingHand","ArmReach","~0.88 m","最大臂展，超出则抓不到"),
        ("ClimbController","gravity","−9.81 m/s²","自由下落加速度")]
tx, ty, tw_ = Inches(0.6), Inches(1.85), Inches(8.6)
tbl = s.shapes.add_table(len(rows), 4, tx, ty, tw_, Inches(3.5)).table
for w_, cw_ in zip(range(4), [Inches(2.4), Inches(2.4), Inches(1.6), Inches(2.2)]):
    tbl.columns[w_].width = cw_
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci); cell.text = val
        para = cell.text_frame.paragraphs[0]; para.runs[0].font.size = Pt(11.5)
        para.runs[0].font.name = "Microsoft YaHei"
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RED
            para.runs[0].font.color.rgb = WHITE; para.runs[0].font.bold = True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else PAPER2
            para.runs[0].font.color.rgb = INK
card(s, Inches(9.4), Inches(1.85), Inches(3.3), Inches(3.5), red_bg=True)
txt(s, Inches(9.65), Inches(2.05), Inches(2.85), Inches(3.2),
    [[("设计取舍", 15, WHITE, True)],
     [("", 6, WHITE, False)],
     [("• 参数按手感设定，非物理推导，已在文档中如实标注；", 12.5, WHITE, False)],
     [("• drain < regen：扶正后能较快回血；", 12.5, WHITE, False)],
     [("• 臂展是真实约束：远处握点须先移动身体才能够到。", 12.5, WHITE, False)]],
    space_after=8, line=1.3)
footer(s, "四 · 关键参数")
print("params done")

# ============== 11 — DEMO ==============
s = slide(); header_band(s)
kicker(s, "五 · 运行演示")
txt(s, Inches(0.55), Inches(0.85), Inches(12.4), Inches(0.7),
    [[("演示：失衡 → 踩脚 → 恢复", 28, INK, True)]])
add_video(s, "demo", Inches(0.6), Inches(1.9), Inches(7.1), Inches(4.0))
txt(s, Inches(0.6), Inches(5.92), Inches(7.1), Inches(0.3),
    [[("演示视频：同侧伸手 → 平衡量下降变红 → 踩橙色脚点 → 恢复", 11, MUT, False)]], align=PP_ALIGN.CENTER)
txt(s, Inches(8.0), Inches(1.9), Inches(4.7), Inches(0.4), [[("握点颜色（用途提示）", 14, RED, True)]])
legend(s, Inches(8.0), Inches(2.45))
bullets(s, Inches(8.0), Inches(2.9), Inches(4.7), Inches(2.6), [
  [("颜色为用途建议，非硬性限制；与真实岩壁一致，手脚可用任意点。", False, INK)],
  [("默认线路设置同侧连续握点：仅靠手会偏出支撑，须踩脚或侧摆才能通过。", False, INK)],
  [("脱落后重生回检查点；登顶显示完成时间。", False, MUT)],
], size=13, gap=10)
footer(s, "五 · 运行演示")
print("demo done")
# <<<NEXT>>>

# ============== 12 — 线路设计 ==============
s = slide(); header_band(s)
kicker(s, "五 · 线路设计")
txt(s, Inches(0.55), Inches(0.82), Inches(12.4), Inches(0.6),
    [[("程序化生成 5 条线路：4 条可完攀 + 1 条用于验证约束", 24, INK, True)]])
routes = [("V1","Warm-up · 入门，熟悉抓握"),("V2","Balance Test · 同侧连续，须用脚"),
          ("V3","The Arête · 易碎点，须快速通过"),("V4","Endurance · 体力管理")]
vw = Inches(2.85); vx = Inches(0.6); vy = Inches(1.65)
for name, cap in routes:
    add_video(s, name, vx, vy, vw, Inches(1.6))
    txt(s, vx, vy+Inches(1.62), vw, Inches(0.35), [[(cap, 10, MUT, True)]], align=PP_ALIGN.CENTER)
    vx = vx + vw + Inches(0.12)
add_video(s, "impossible", Inches(0.6), Inches(3.95), Inches(3.6), Inches(2.0))
txt(s, Inches(0.6), Inches(5.97), Inches(3.6), Inches(0.3),
    [[("The Gap · 中段约 2m 空白（不可完攀）", 10, RED, True)]], align=PP_ALIGN.CENTER)
bullets(s, Inches(4.5), Inches(4.0), Inches(8.2), Inches(2.0), [
  [("The Gap 用于验证臂展约束：", True, REDD), ("中段留约 2m 空白，手臂伸直仍够不到，说明臂展限制确实生效，而非装饰。", False, INK)],
  [("生成方式：", True, REDD), ("RouteBuilder 由基本几何体程序化构建墙体、彩色握点与登顶触发区，无需美术资源即可运行。", False, INK)],
  [("难度递增设计使「脚法是否被使用」成为可观测的行为差异，便于后续评测。", False, MUT)],
], size=13.5, gap=11)
footer(s, "五 · 线路设计")
print("routes done")

# ============== 13 — 开发与验证方法 ==============
s = slide(); header_band(s)
kicker(s, "六 · 开发与验证方法")
txt(s, Inches(0.55), Inches(0.85), Inches(12.4), Inches(0.7),
    [[("以自动化验证驱动开发", 28, INK, True)]])
metrics = [("10/10","端到端自检","脚本驱动真实玩法：失衡→脱落→重生→登顶全流程通过"),
           ("9/9","算法单元测试","平衡与位移核心数学（ClimbMath）逐项通过"),
           ("2","可玩验证视角","第三人称与第一人称场景，键鼠即可试玩")]
cx = Inches(0.6)
for big, h, body in metrics:
    card(s, cx, Inches(1.95), Inches(4.0), Inches(2.0))
    txt(s, cx, Inches(2.1), Inches(4.0), Inches(0.7), [[(big, 38, RED, True)]], align=PP_ALIGN.CENTER)
    txt(s, cx, Inches(2.9), Inches(4.0), Inches(0.35), [[(h, 14, INK, True)]], align=PP_ALIGN.CENTER)
    txt(s, cx+Inches(0.25), Inches(3.3), Inches(3.5), Inches(0.65), [[(body, 10.5, MUT, False)]], align=PP_ALIGN.CENTER, line=1.2)
    cx = cx + Inches(4.05)
c = card(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(1.15), red_bg=True)
txt(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.95),
    [[("方法说明", 15, WHITE, True)],
     [("先用脚本「虚拟玩家」对真实玩法做端到端自检，再以无头渲染产出演示并逐帧比对打磨；每次改动均回归自检。机制正确性由自动化保证，五名成员得以并行开发，不依赖共享头显。", 13, WHITE, False)]],
    space_after=4, line=1.3)
footer(s, "六 · 开发与验证方法")
print("verify done")

# ============== 14 — 评测方案 ==============
s = slide(); header_band(s)
kicker(s, "六 · 评测方案")
txt(s, Inches(0.55), Inches(0.85), Inches(12.4), Inches(0.7),
    [[("评测方案已设计，", 28, INK, True), ("待开展真人实验", 28, RED, True)]])
bullets(s, Inches(0.6), Inches(1.9), Inches(6.4), Inches(2.2), [
  [("实验设计：", True, REDD), ("被试内对比，A 纯手版 vs B 平衡+脚法版，同线路、顺序平衡。", False, INK)],
  [("客观指标：", True, REDD), ("登顶时间、按原因分类的跌落次数（松手 / 失衡 / 体力）、平衡余量曲线。", False, INK)],
  [("主观量表：", True, REDD), ("NASA-TLX 工作负荷、SSQ 模拟器晕动、真实感与偏好量表。", False, INK)],
], size=13.5, gap=10)
c = card(s, Inches(0.6), Inches(4.0), Inches(6.4), Inches(1.4), red_bg=True)
txt(s, Inches(0.85), Inches(4.15), Inches(5.9), Inches(1.1),
    [[("研究假设（待检验）", 14, WHITE, True)],
     [("相比 A，B 的真实感与挑战感更高，而模拟器晕动不显著上升——难度来自决策判断而非画面运动。", 13, WHITE, False)]],
    space_after=4, line=1.3)
txt(s, Inches(7.3), Inches(1.9), Inches(5.4), Inches(0.4), [[("已完成 ✓ / 待完成 ○", 14, RED, True)]])
bullets(s, Inches(7.3), Inches(2.4), Inches(5.4), Inches(3.0), [
  [("✓ 平衡机制与脚法建模", False, INK)],
  [("✓ 5 条线路 + 自动化自检（10/10 · 9/9）", False, INK)],
  [("✓ HUD、音效、双视角可玩", False, INK)],
  [("✓ 演示视频、报告初稿、评测协议与采集工具", False, INK)],
  [("○ 招募 5–8 名被试，开展实验并统计数据", False, RED)],
  [("○ 真机构建与帧率 / 舒适度优化（可选）", False, MUT)],
], size=13, gap=8)
footer(s, "六 · 评测方案")
print("eval done")

# ============== 15 — 总结 ==============
s = slide()
rect(s, 0, 0, Inches(7.6), SH, WHITE)
rect(s, Inches(7.6), 0, SW-Inches(7.6), SH, RED)
grad_band(s, 0, 0, Inches(7.6), Inches(0.16))
txt(s, Inches(0.6), Inches(1.1), Inches(6.6), Inches(0.4), [[("六 · 总结与展望", 14, RED, True)]])
txt(s, Inches(0.55), Inches(1.6), Inches(6.8), Inches(1.0), [[("总结", 38, INK, True)]])
bullets(s, Inches(0.6), Inches(2.7), Inches(6.7), Inches(2.6), [
  [("提出并实现了一种在三点追踪下的平衡与脚法机制：", True, REDD)],
  [("• 以头显近似重心、以横向支撑区间判定失衡；", False, INK)],
  [("• 以自动吸附的虚拟脚提供支撑，使脚法影响成败；", False, INK)],
  [("• 系统完整可运行，并通过自动化端到端验证。", False, INK)],
  [("展望：", True, REDD), ("开展真人评测，并向二维支撑、玩家自主选点、真机部署扩展。", False, INK)],
], size=13, gap=8)
txt(s, Inches(0.6), Inches(6.7), Inches(6.8), Inches(0.4),
    [[("项目代码：github.com/Icecream0507/ai3618-vr-climbing", 11, MUT, False)]])
# right panel
txt(s, Inches(8.0), Inches(1.5), Inches(4.8), Inches(0.6), [[("核心贡献", 22, WHITE, True)]])
rect(s, Inches(8.0), Inches(2.15), Inches(4.5), Pt(2), RGBColor(0xFF,0xB0,0xBC))
txt(s, Inches(8.0), Inches(2.45), Inches(4.8), Inches(2.6),
    [[("在不增加硬件的前提下，", 15, WHITE, False)],
     [("用头显近似重心、用虚拟脚提供支撑，", 15, RGBColor(0xFF,0xD9,0xDF), True)],
     [("将真实攀岩的平衡与脚法，", 15, WHITE, False)],
     [("重建为消费级 VR 上可玩、可验证的机制。", 15, WHITE, False)]],
    space_after=6, line=1.4)
txt(s, Inches(8.0), Inches(5.7), Inches(4.8), Inches(0.8), [[("谢谢 · 欢迎提问", 30, WHITE, True)]])
print("summary done")

out = os.path.join(HERE, "SummitVR_Pre.pptx")
prs.save(out)
print("SAVED ->", out)


